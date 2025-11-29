<!-- filepath: /home/abdul-rehman/Desktop/CS/Web-Technologies/Mid assignment/convert_to_ppt.py -->
#!/usr/bin/env python3
"""
HTML to PowerPoint Converter for PHP Presentation
Requires: pip install python-pptx beautifulsoup4 lxml
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup
import re

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def add_title_slide(prs, soup_slide):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(102, 126, 234)
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2), Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = soup_slide.find('h1').text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    h2 = soup_slide.find('h2')
    if h2:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.5), Inches(9), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = h2.text
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(32)
        subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
        subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Tech badges
    badges = soup_slide.find_all('span', class_='badge')
    if badges:
        badge_text = " • ".join([badge.text for badge in badges])
        badge_box = slide.shapes.add_textbox(
            Inches(1), Inches(5), Inches(8), Inches(0.5)
        )
        badge_frame = badge_box.text_frame
        badge_frame.text = badge_text
        badge_para = badge_frame.paragraphs[0]
        badge_para.font.size = Pt(18)
        badge_para.font.color.rgb = RGBColor(255, 255, 255)
        badge_para.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, soup_slide, slide_num):
    """Create content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add slide number
    slide_num_box = slide.shapes.add_textbox(
        Inches(9), Inches(7), Inches(0.5), Inches(0.3)
    )
    slide_num_frame = slide_num_box.text_frame
    slide_num_frame.text = str(slide_num)
    slide_num_para = slide_num_frame.paragraphs[0]
    slide_num_para.font.size = Pt(12)
    slide_num_para.font.color.rgb = RGBColor(102, 126, 234)
    
    # Title
    h1 = soup_slide.find('h1')
    if h1:
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = h1.text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(102, 126, 234)
    
    current_y = 1.3
    
    # Process content elements
    for element in soup_slide.find_all(['h2', 'h3', 'p', 'ul', 'div', 'table']):
        if element.name == 'h2':
            text_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(current_y), Inches(9), Inches(0.5)
            )
            text_frame = text_box.text_frame
            text_frame.text = element.text
            para = text_frame.paragraphs[0]
            para.font.size = Pt(24)
            para.font.bold = True
            para.font.color.rgb = RGBColor(118, 75, 162)
            current_y += 0.6
            
        elif element.name == 'h3':
            text_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(current_y), Inches(9), Inches(0.4)
            )
            text_frame = text_box.text_frame
            text_frame.text = element.text
            para = text_frame.paragraphs[0]
            para.font.size = Pt(20)
            para.font.bold = True
            para.font.color.rgb = RGBColor(102, 126, 234)
            current_y += 0.5
            
        elif element.name == 'p' and 'highlight-box' not in element.get('class', []):
            text_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(current_y), Inches(9), Inches(0.4)
            )
            text_frame = text_box.text_frame
            text_frame.text = element.text
            para = text_frame.paragraphs[0]
            para.font.size = Pt(16)
            para.font.color.rgb = RGBColor(51, 51, 51)
            current_y += 0.5
            
        elif element.name == 'ul':
            for li in element.find_all('li'):
                if current_y > 6.5:
                    break
                text_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(current_y), Inches(8.5), Inches(0.35)
                )
                text_frame = text_box.text_frame
                text_frame.text = "• " + li.text
                para = text_frame.paragraphs[0]
                para.font.size = Pt(16)
                para.font.color.rgb = RGBColor(51, 51, 51)
                current_y += 0.4
            current_y += 0.2
            
        elif element.name == 'div' and 'code-block' in element.get('class', []):
            if current_y > 6:
                continue
            # Clean code text
            code_text = element.get_text()
            code_text = re.sub(r'\s+', ' ', code_text).strip()
            
            code_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(current_y), Inches(9), Inches(1.5)
            )
            code_frame = code_box.text_frame
            code_frame.text = code_text[:300]  # Limit length
            code_frame.word_wrap = True
            para = code_frame.paragraphs[0]
            para.font.name = 'Consolas'
            para.font.size = Pt(12)
            para.font.color.rgb = RGBColor(212, 212, 212)
            
            # Background
            shape = code_box
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(30, 30, 30)
            
            current_y += 1.7
            
        elif element.name == 'div' and 'highlight-box' in element.get('class', []):
            if current_y > 6:
                continue
            highlight_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(current_y), Inches(9), Inches(0.8)
            )
            highlight_frame = highlight_box.text_frame
            highlight_frame.text = element.get_text()
            para = highlight_frame.paragraphs[0]
            para.font.size = Pt(14)
            para.font.color.rgb = RGBColor(102, 126, 234)
            
            # Background
            fill = highlight_box.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(240, 242, 255)
            
            current_y += 1
            
        elif element.name == 'table':
            # Add table (simplified)
            rows = element.find_all('tr')
            if rows and current_y < 5:
                table_text = []
                for row in rows[:5]:  # Limit rows
                    cells = row.find_all(['th', 'td'])
                    row_text = " | ".join([cell.text for cell in cells])
                    table_text.append(row_text)
                
                table_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(current_y), Inches(9), Inches(2)
                )
                table_frame = table_box.text_frame
                table_frame.text = "\n".join(table_text)
                para = table_frame.paragraphs[0]
                para.font.name = 'Consolas'
                para.font.size = Pt(12)
                
                current_y += 2.2

def create_presentation():
    """Main function to create PowerPoint presentation"""
    print("📊 Creating PowerPoint presentation...")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Read HTML
    print("📖 Reading HTML file...")
    try:
        with open('presentation_clean.html', 'r', encoding='utf-8') as f:
            html_content = f.read()
    except FileNotFoundError:
        print("❌ Error: presentation_clean.html not found!")
        print("Please make sure the HTML file exists in the current directory.")
        return
    
    soup = BeautifulSoup(html_content, 'html.parser')
    slides_html = soup.find_all('div', class_='slide')
    
    print(f"✅ Found {len(slides_html)} slides")
    
    # Process each slide
    for idx, slide_html in enumerate(slides_html, 1):
        print(f"Processing slide {idx}/{len(slides_html)}...")
        
        # Check if title slide
        is_title = 'title-slide' in slide_html.get('class', [])
        
        if is_title:
            add_title_slide(prs, slide_html)
        else:
            add_content_slide(prs, slide_html, idx)
    
    # Save presentation
    output_file = 'PHP_CRUD_Presentation.pptx'
    print(f"💾 Saving presentation as {output_file}...")
    prs.save(output_file)
    print(f"✅ Success! Presentation saved as {output_file}")
    print(f"📍 Total slides created: {len(slides_html)}")

if __name__ == '__main__':
    print("=" * 50)
    print("HTML to PowerPoint Converter")
    print("=" * 50)
    create_presentation()